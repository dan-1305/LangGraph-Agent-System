import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_dashboard_ppt():
    prs = Presentation()

    # Define some colors
    BLUE = RGBColor(0, 112, 192)
    GREEN = RGBColor(0, 176, 80)
    DARK_GRAY = RGBColor(89, 89, 89)

    # 1. Slide Bìa (Title Slide)
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "PHÂN TÍCH VÀ DỰ BÁO\nGIÁ BẤT ĐỘNG SẢN TP.HCM & ĐỒNG NAI"
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = BLUE
    
    subtitle.text = "Ứng dụng Machine Learning & End-to-End Pipeline\n\nBáo cáo Tiểu Luận"

    # 2. Slide Tổng quan (Dashboard Style)
    slide_layout = prs.slide_layouts[5] # Blank slide with title
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "PROJECT DASHBOARD OVERVIEW"
    slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = BLUE

    # KPI 1: Số dòng dữ liệu
    txbox1 = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(3), Inches(1.5))
    tf1 = txbox1.text_frame
    p1 = tf1.paragraphs[0]
    p1.text = "914"
    p1.font.size = Pt(60)
    p1.font.bold = True
    p1.font.color.rgb = BLUE
    p1.alignment = PP_ALIGN.CENTER
    p1_sub = tf1.add_paragraph()
    p1_sub.text = "Bất động sản hợp lệ"
    p1_sub.font.size = Pt(20)
    p1_sub.alignment = PP_ALIGN.CENTER

    # KPI 2: Features
    txbox2 = slide.shapes.add_textbox(Inches(5.5), Inches(2), Inches(3), Inches(1.5))
    tf2 = txbox2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = "11"
    p2.font.size = Pt(60)
    p2.font.bold = True
    p2.font.color.rgb = GREEN
    p2.alignment = PP_ALIGN.CENTER
    p2_sub = tf2.add_paragraph()
    p2_sub.text = "Features (Đã Regex)"
    p2_sub.font.size = Pt(20)
    p2_sub.alignment = PP_ALIGN.CENTER

    # KPI 3: R2 High End
    txbox3 = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(3), Inches(1.5))
    tf3 = txbox3.text_frame
    p3 = tf3.paragraphs[0]
    p3.text = "0.45"
    p3.font.size = Pt(60)
    p3.font.bold = True
    p3.font.color.rgb = DARK_GRAY
    p3.alignment = PP_ALIGN.CENTER
    p3_sub = tf3.add_paragraph()
    p3_sub.text = "R² High-End (MAE: 4.6 Tỷ)"
    p3_sub.font.size = Pt(16)
    p3_sub.alignment = PP_ALIGN.CENTER

    # KPI 4: R2 Mass Market
    txbox4 = slide.shapes.add_textbox(Inches(5.5), Inches(4.5), Inches(3), Inches(1.5))
    tf4 = txbox4.text_frame
    p4 = tf4.paragraphs[0]
    p4.text = "0.37"
    p4.font.size = Pt(60)
    p4.font.bold = True
    p4.font.color.rgb = DARK_GRAY
    p4.alignment = PP_ALIGN.CENTER
    p4_sub = tf4.add_paragraph()
    p4_sub.text = "R² Mass-Market (MAE: 1.9 Tỷ)"
    p4_sub.font.size = Pt(16)
    p4_sub.alignment = PP_ALIGN.CENTER

    # 3. Slide Feature Engineering
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "FEATURE ENGINEERING (KỸ THUẬT LÕI)"
    tf = slide.placeholders[1].text_frame
    tf.text = "Biến đổi văn bản thành tín hiệu toán học:"
    
    p = tf.add_paragraph()
    p.text = "Regex Binarization: Rút trích 'is_frontage', 'is_alley', 'has_furniture' từ Tiêu đề."
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Xử lý Outlier: Loại bỏ nhiễu bằng Z-Score cục bộ theo từng Quận/Huyện."
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Log Transform: Chuẩn hóa phân phối đuôi dài của Giá (Price_VND)."
    p.level = 1

    # 4. Slide Kiến trúc phân mảnh
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "KIẾN TRÚC PHÂN MẢNH (SEGMENTATION)"
    tf = slide.placeholders[1].text_frame
    tf.text = "Giải quyết tính Không đồng nhất (Heterogeneity) của BĐS:"
    
    p = tf.add_paragraph()
    p.text = "Phân khúc Cao cấp (High-End): Biệt thự, Nhà mặt tiền, Quận trung tâm."
    p.level = 1
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "Phân khúc Phổ thông (Mass-Market): Các trường hợp còn lại."
    p.level = 1
    p.font.bold = True

    p = tf.add_paragraph()
    p.text = "Mỗi phân khúc sử dụng một mô hình XGBoost riêng biệt để tăng độ chính xác."
    p.level = 0

    # 5. Slide Tối ưu hóa
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "TỐI ƯU HÓA & TRÁNH OVERFITTING"
    tf = slide.placeholders[1].text_frame
    tf.text = "Pipeline và Hyperparameter Tuning:"
    
    p = tf.add_paragraph()
    p.text = "Sử dụng Scikit-Learn Pipeline (StandardScaler, OneHotEncoder) để ngăn Data Leakage."
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Áp dụng RandomizedSearchCV với K-Fold Cross Validation (cv=3)."
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Tinh chỉnh: Giảm max_depth, tăng reg_alpha & reg_lambda để phạt mô hình."
    p.level = 1

    # 6. Slide Demo
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "TRIỂN KHAI & DEMO (DEPLOYMENT)"
    tf = slide.placeholders[1].text_frame
    tf.text = "Mang mô hình vào thực tiễn:"
    
    p = tf.add_paragraph()
    p.text = "Đóng gói mô hình (.pkl) thành Web App sử dụng Flask."
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Tích hợp bản đồ nhiệt (Heatmap) bằng Folium trực quan."
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Hệ thống tự động điều hướng input người dùng vào Model tương ứng (High-End/Mass)."
    p.level = 1

    # Save
    output_path = os.path.join("projects", "real_estate_prediction", "reports", "Bao_Cao_Tieu_Luan_BDS.pptx")
    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    prs.save(output_path)
    print(f"Đã tạo thành công file PowerPoint tại: {output_path}")

if __name__ == "__main__":
    create_dashboard_ppt()
